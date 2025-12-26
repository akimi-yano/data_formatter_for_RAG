
import os
from pypdf import PdfReader
from docx import Document
from pptx import Presentation

def extract_text_from_pdf(file_path: str) -> str:
    """Extracts text from a PDF file."""
    try:
        reader = PdfReader(file_path)
        text = ""
        for page in reader.pages:
            text += page.extract_text() + "\n"
        return text
    except Exception as e:
        return f"Error extracting PDF: {str(e)}"

def extract_text_from_docx(file_path: str) -> str:
    """Extracts text from a DOCX file."""
    try:
        doc = Document(file_path)
        text = "\n".join([para.text for para in doc.paragraphs])
        return text
    except Exception as e:
        return f"Error extracting DOCX: {str(e)}"

def extract_text_from_pptx(file_path: str) -> str:
    """Extracts text from a PPTX file."""
    try:
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        return f"Error extracting PPTX: {str(e)}"

def extract_text(file_path: str) -> str:
    """
    Detects file extension and extracts text accordingly.
    Supported formats: .pdf, .docx, .pptx
    """
    ext = os.path.splitext(file_path)[1].lower()
    
    if ext == ".pdf":
        return extract_text_from_pdf(file_path)
    elif ext == ".docx":
        return extract_text_from_docx(file_path)
    elif ext == ".pptx":
        return extract_text_from_pptx(file_path)
    else:
        return f"Unsupported file format: {ext}"
